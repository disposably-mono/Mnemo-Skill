"""Source ingestion: study material -> normalized text chunks with provenance.

This is the Anki-agnostic front of the pipeline. It turns a file into uniform
``Chunk``s (``text`` + a human-readable ``source`` like ``lecture.pdf p.4``) that
Claude then reads to author Facts. Markdown/plain text is read directly; PDFs go
through PyMuPDF (one chunk per page); slide decks through python-pptx (one chunk
per slide). SKILL.md step 2 runs:

    python scripts/ingest.py <file>
"""

from __future__ import annotations

import re
import sys
from dataclasses import dataclass
from pathlib import Path

if __package__ in (None, ""):  # allow `python scripts/ingest.py <file>`
    sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

_TEXT_EXTS = {".md", ".markdown", ".txt", ".text"}
MIN_IMAGE_PX = 32
_GREEK = re.compile(r"[\u0370-\u03ff\u1f00-\u1fff]")
_MATH_OPERATOR = re.compile(
    r"(?:[=+*/^<>\u00b1\u00d7\u00f7\u2192\u2194\u2211\u220f\u221a\u221d\u221e\u222b\u2248\u2260\u2264\u2265\u2202\u2207]|->|=>|<=|>=)"
)
_DISTINCTIVE_MATH = re.compile(
    r"[\u00b1\u00d7\u00f7\u2192\u2194\u2211\u220f\u221a\u221d\u221e\u222b\u2248\u2260\u2264\u2265\u2202\u2207]"
)


@dataclass
class Chunk:
    """A normalized unit of source text with its provenance."""

    text: str
    source: str


def ingest(
    path: str | Path,
    *,
    ocr: bool = False,
    extract_images: Path | None = None,
) -> list[Chunk]:
    """Parse a source file into normalized text chunks. Dispatch by extension.

    ``ocr`` only affects PDFs: when true, image-only pages are sent through
    PyMuPDF's Tesseract integration before falling back to a visible marker.
    ``extract_images`` opts PDFs into saving embedded raster source visuals.
    """
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(path)

    ext = path.suffix.lower()
    if ext in _TEXT_EXTS:
        return _ingest_text(path)
    if ext == ".pdf":
        return _ingest_pdf(path, ocr=ocr, extract_images=extract_images)
    if ext == ".pptx":
        return _ingest_pptx(path)
    raise ValueError(f"unsupported file type: {path.suffix!r} ({path.name})")


def _ingest_text(path: Path) -> list[Chunk]:
    text = path.read_text(encoding="utf-8").strip()
    if not text:
        return []
    return [Chunk(text=text, source=path.name)]


def _scanned_page_marker(image_count: int, ocr_attempted: bool) -> str:
    """Visible placeholder for a page that carries images but no text layer.

    Surfacing the gap keeps the missing source material auditable instead of
    dropping it silently, per the invariant that unsupported material stays
    visible. The reader (human or Claude) sees a page needs OCR or manual
    transcription before any card can be grounded in it.
    """
    plural = "image" if image_count == 1 else "images"
    remedy = (
        "OCR produced no text; transcribe manually"
        if ocr_attempted
        else "enable --ocr or transcribe manually"
    )
    return (
        f"[image-only page: no extractable text layer; {image_count} embedded "
        f"{plural}. Likely scanned or a figure. {remedy} before grounding cards "
        "here.]"
    )


def _ocr_page_text(page) -> str:
    """Best-effort OCR of a single page; empty string when OCR is unavailable.

    PyMuPDF delegates to a Tesseract install. A missing binary or tessdata
    raises, which we treat as "no OCR available" rather than failing ingestion.
    OCR output is lower confidence, so callers mark its provenance as ``(OCR)``.
    """
    try:
        textpage = page.get_textpage_ocr(full=True)
        return page.get_text(textpage=textpage).strip()
    except Exception:  # missing tesseract/tessdata, unsupported page, etc.
        return ""


def _format_pdf_table(rows: list[list[str | None]]) -> str:
    """Render extracted table rows as pipe-joined lines (the pptx convention)."""
    lines: list[str] = []
    for row in rows:
        line = " | ".join((cell or "").strip() for cell in row)
        if line.strip(" |"):  # drop fully empty rows
            lines.append(line)
    return "\n".join(lines)


def _extract_pdf_tables(page) -> list[str]:
    """Structured pipe-row rendering of any tables PyMuPDF detects on a page.

    Flat ``page.get_text()`` reads a table column-blind, scrambling which value
    belongs to which row/header. Re-emitting the detected grid preserves those
    relationships so comparison and lookup cards stay grounded. Conservative:
    only grids with at least two rows and two columns are kept, to avoid
    mislabeling ordinary prose as tabular.
    """
    try:
        finder = page.find_tables()
    except Exception:  # detection is best-effort; never fail ingestion over it
        return []
    blocks: list[str] = []
    for table in finder.tables:
        rows = table.extract()
        if len(rows) < 2 or max((len(row) for row in rows), default=0) < 2:
            continue
        rendered = _format_pdf_table(rows)
        if rendered:
            blocks.append(f"Table:\n{rendered}")
    return blocks


def _looks_like_pdf_math(text: str) -> bool:
    """Return true only for short text-layer lines with strong math signals."""
    value = " ".join(text.split())
    if not value or len(value) > 240:
        return False
    operators = _MATH_OPERATOR.findall(value)
    if not operators:
        return False
    if _DISTINCTIVE_MATH.search(value):
        return True
    if _GREEK.search(value) and operators:
        return True
    if "=" in value and re.search(r"[A-Za-z0-9]", value):
        left, _, right = value.partition("=")
        return bool(left.strip() and right.strip())
    return len(operators) >= 2 and bool(re.search(r"[A-Za-z0-9]", value))


def _extract_pdf_math(page, source: str) -> list[str]:
    """Label likely equations while preserving text-layer span characters.

    PyMuPDF does not parse mathematical semantics, but its structured text keeps
    font-split spans on the same line. Joining those spans lets Greek letters and
    operators survive as an auditable candidate without claiming OCR or LaTeX
    reconstruction accuracy.
    """
    try:
        blocks = page.get_text("dict", sort=True).get("blocks", [])
    except Exception:
        return []
    markers: list[str] = []
    seen: set[str] = set()
    for block in blocks:
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            text = "".join(
                str(span.get("text", "")) for span in line.get("spans", [])
            ).strip()
            normalized = " ".join(text.split())
            if not _looks_like_pdf_math(normalized) or normalized in seen:
                continue
            seen.add(normalized)
            markers.append(f"[math: {normalized} | {source} text layer]")
    return markers


def _is_qualifying_image(info: dict) -> bool:
    """Heuristically reject tiny embedded assets that are usually decoration."""
    return (
        info.get("width", 0) >= MIN_IMAGE_PX
        and info.get("height", 0) >= MIN_IMAGE_PX
    )


def _extract_pdf_images(
    doc, page, path: Path, page_number: int, output: Path
) -> list[str]:
    """Save qualifying page images and return auditable figure marker lines.

    Image extraction is best-effort because malformed embedded assets should
    not prevent the page's text from being ingested.
    """
    try:
        images = page.get_images(full=True)
    except Exception:
        return []
    markers: list[str] = []
    seen_xrefs: set[int] = set()
    for image in images:
        xref = image[0]
        if xref in seen_xrefs:
            continue
        seen_xrefs.add(xref)
        try:
            info = doc.extract_image(xref)
            if not _is_qualifying_image(info):
                continue
            image_number = len(markers) + 1
            filename = (
                f"{path.stem}-p{page_number}-img{image_number}.{info['ext']}"
            )
            saved = output / filename
            saved.write_bytes(info["image"])
        except Exception:
            continue
        markers.append(f"[figure: {saved} | {path.name} p.{page_number}]")
    return markers


def _ingest_pdf(
    path: Path,
    *,
    ocr: bool = False,
    extract_images: Path | None = None,
) -> list[Chunk]:
    import fitz  # PyMuPDF

    chunks: list[Chunk] = []
    if extract_images is not None:
        extract_images.mkdir(parents=True, exist_ok=True)
    with fitz.open(str(path)) as doc:
        for index, page in enumerate(doc, start=1):
            source = f"{path.name} p.{index}"
            figures = (
                _extract_pdf_images(doc, page, path, index, extract_images)
                if extract_images is not None
                else []
            )
            text = page.get_text().strip()
            if text:
                parts = [
                    text,
                    *_extract_pdf_tables(page),
                    *_extract_pdf_math(page, source),
                    *figures,
                ]
                chunks.append(Chunk(text="\n".join(parts), source=source))
                continue
            image_count = len(page.get_images(full=True))
            if image_count == 0:
                continue  # genuinely blank page -> nothing to surface
            ocr_text = _ocr_page_text(page) if ocr else ""
            if ocr_text:
                parts = [ocr_text, *figures]
                chunks.append(Chunk(text="\n".join(parts), source=f"{source} (OCR)"))
            else:
                marker = _scanned_page_marker(image_count, ocr_attempted=ocr)
                chunks.append(Chunk(text="\n".join([marker, *figures]), source=source))
    return chunks


def _format_chart_value(value) -> str:
    """Render chart values without unnecessary floating-point noise."""
    if value is None:
        return ""
    if isinstance(value, float):
        if value.is_integer():
            return str(int(value))
        return format(value, ".12g")
    return str(value).strip()


def _extract_pptx_chart(chart) -> str:
    """Render accessible chart categories and series as a relational table.

    This preserves the values stored in the presentation. It intentionally does
    not infer trends, causes, or conclusions that are not explicit source text.
    """
    title = ""
    try:
        if chart.has_title:
            title = chart.chart_title.text_frame.text.strip()
    except Exception:
        pass

    series = list(chart.series)
    names = [
        str(item.name).strip() or f"Series {index}"
        for index, item in enumerate(series, start=1)
    ]
    try:
        categories = [
            str(getattr(category, "label", category)).strip()
            for category in chart.plots[0].categories
        ]
    except Exception:
        categories = []

    lines = ["Chart:"]
    if title:
        lines.append(f"Title: {title}")
    if categories and series:
        lines.append("Category | " + " | ".join(names))
        values = [list(item.values) for item in series]
        for index, category in enumerate(categories):
            row = [
                _format_chart_value(series_values[index])
                if index < len(series_values)
                else ""
                for series_values in values
            ]
            lines.append(f"{category} | " + " | ".join(row))
    elif names:
        lines.append("Series: " + " | ".join(names))
    return "\n".join(lines) if len(lines) > 1 else ""


def _ingest_pptx(path: Path) -> list[Chunk]:
    from pptx import Presentation

    chunks: list[Chunk] = []
    prs = Presentation(str(path))
    for index, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            elif shape.has_table:
                rows = [
                    " | ".join(cell.text.strip() for cell in row.cells)
                    for row in shape.table.rows
                ]
                table = "\n".join(row for row in rows if row.strip(" |"))
                if table:
                    parts.append(table)
            elif shape.has_chart:
                chart = _extract_pptx_chart(shape.chart)
                if chart:
                    parts.append(chart)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Speaker notes:\n{notes}")
        text = "\n".join(parts).strip()
        if text:  # skip empty slides
            chunks.append(Chunk(text=text, source=f"{path.name} slide {index}"))
    return chunks


def main(argv: list[str] | None = None) -> int:
    import argparse

    parser = argparse.ArgumentParser(
        description="Normalize study material into text chunks with provenance."
    )
    parser.add_argument("file", help="Path to a .pdf, .pptx, .md, or .txt source.")
    parser.add_argument(
        "--ocr",
        action="store_true",
        help="OCR image-only PDF pages (requires Tesseract); falls back to a "
        "visible marker when OCR is unavailable.",
    )
    parser.add_argument(
        "--extract-images",
        type=Path,
        metavar="DIR",
        help="Save qualifying embedded PDF images and add figure provenance.",
    )
    args = parser.parse_args(argv)

    for chunk in ingest(args.file, ocr=args.ocr, extract_images=args.extract_images):
        print(f"--- {chunk.source} ---")
        print(chunk.text)
        print()
    return 0


if __name__ == "__main__":  # pragma: no cover
    sys.exit(main())
