"""Tests for source ingestion (scripts/ingest.py).

Fixtures are built programmatically with the real libraries (PyMuPDF for PDF,
python-pptx for slides) rather than mocked, so these exercise the actual parse
paths. Markdown/text is read directly.
"""

import subprocess
import sys
from pathlib import Path

import pytest

from scripts.ingest import Chunk, _extract_pdf_math, ingest

_REPO_ROOT = Path(__file__).resolve().parent.parent


def test_markdown_is_read_as_a_single_chunk(tmp_path):
    md = tmp_path / "notes.md"
    md.write_text("# Title\n\nThe mitochondria is the powerhouse of the cell.\n")

    chunks = ingest(md)

    assert len(chunks) == 1
    assert isinstance(chunks[0], Chunk)
    assert "powerhouse" in chunks[0].text
    assert chunks[0].source == "notes.md"


def test_plain_text_is_supported(tmp_path):
    txt = tmp_path / "scratch.txt"
    txt.write_text("just some pasted notes")
    chunks = ingest(txt)
    assert chunks[0].text == "just some pasted notes"
    assert chunks[0].source == "scratch.txt"


def test_pdf_yields_one_chunk_per_nonempty_page(tmp_path):
    import fitz  # PyMuPDF

    pdf = tmp_path / "lecture.pdf"
    doc = fitz.open()
    doc.new_page().insert_text((72, 72), "Photosynthesis converts light.")
    doc.new_page()  # intentionally blank -> should be skipped
    doc.new_page().insert_text((72, 72), "ATP is energy currency.")
    doc.save(str(pdf))
    doc.close()

    chunks = ingest(pdf)

    assert len(chunks) == 2  # blank page dropped
    assert "Photosynthesis" in chunks[0].text
    assert chunks[0].source == "lecture.pdf p.1"
    assert "ATP" in chunks[1].text
    assert chunks[1].source == "lecture.pdf p.3"  # provenance keeps real page no.


def test_pptx_yields_one_chunk_per_slide_with_all_shape_text(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    pptx = tmp_path / "deck.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.text_frame.text = "Newton's first law: inertia."
    prs.save(str(pptx))

    chunks = ingest(pptx)

    assert len(chunks) == 1
    assert "inertia" in chunks[0].text
    assert chunks[0].source == "deck.pptx slide 1"


def test_pptx_includes_tables_and_existing_speaker_notes(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    pptx = tmp_path / "structured.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(5), Inches(2)).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Meaning"
    table.cell(1, 0).text = "ROI"
    table.cell(1, 1).text = "Return on investment"
    slide.notes_slide.notes_text_frame.text = "Explain why cost belongs in the denominator."
    prs.save(str(pptx))

    chunks = ingest(pptx)

    assert "Metric | Meaning" in chunks[0].text
    assert "ROI | Return on investment" in chunks[0].text
    assert "Speaker notes:" in chunks[0].text
    assert "cost belongs in the denominator" in chunks[0].text


def _image_only_pdf(path: Path, *, image_px: int = 40, text: str | None = None) -> None:
    """Write a one-page PDF with an embedded image and optional text layer."""
    import fitz  # PyMuPDF

    doc = fitz.open()
    page = doc.new_page()
    if text:
        page.insert_text((72, 72), text)
    pix = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, image_px, image_px))
    pix.clear_with(200)
    page.insert_image(fitz.Rect(10, 10, 10 + image_px, 10 + image_px), pixmap=pix)
    doc.save(str(path))
    doc.close()


def test_image_only_pdf_page_is_surfaced_not_dropped(tmp_path):
    pdf = tmp_path / "scanned.pdf"
    _image_only_pdf(pdf)

    chunks = ingest(pdf)

    # The page carries no text layer but must remain visible, not silently lost.
    assert len(chunks) == 1
    assert chunks[0].source == "scanned.pdf p.1"
    assert "image-only page" in chunks[0].text
    assert "1 embedded image" in chunks[0].text
    assert "--ocr" in chunks[0].text  # tells the reader how to recover it


def test_truly_blank_pdf_page_is_still_skipped(tmp_path):
    import fitz  # PyMuPDF

    pdf = tmp_path / "mixed.pdf"
    doc = fitz.open()
    doc.new_page().insert_text((72, 72), "Real content here.")
    doc.new_page()  # blank: no text, no images -> nothing to surface
    doc.save(str(pdf))
    doc.close()

    chunks = ingest(pdf)

    assert len(chunks) == 1
    assert "Real content" in chunks[0].text
    assert chunks[0].source == "mixed.pdf p.1"


def test_ocr_flag_falls_back_to_marker_when_ocr_unavailable(tmp_path):
    # No Tesseract in this environment: --ocr must degrade gracefully to the
    # marker rather than raising or dropping the page.
    pdf = tmp_path / "scanned.pdf"
    _image_only_pdf(pdf)

    chunks = ingest(pdf, ocr=True)

    assert len(chunks) == 1
    assert chunks[0].source == "scanned.pdf p.1"  # no "(OCR)" suffix on fallback
    assert "image-only page" in chunks[0].text
    assert "OCR produced no text" in chunks[0].text


def test_ocr_recovered_text_is_labeled_in_provenance(tmp_path, monkeypatch):
    # When OCR succeeds, the recovered text is emitted with an (OCR) provenance
    # tag so downstream grounding can treat it as lower confidence.
    import scripts.ingest as ingest_mod

    pdf = tmp_path / "scanned.pdf"
    _image_only_pdf(pdf)
    monkeypatch.setattr(ingest_mod, "_ocr_page_text", lambda page: "Recovered text.")

    chunks = ingest(pdf, ocr=True)

    assert len(chunks) == 1
    assert chunks[0].text == "Recovered text."
    assert chunks[0].source == "scanned.pdf p.1 (OCR)"


def test_extract_images_writes_file_and_adds_figure_marker(tmp_path):
    pdf = tmp_path / "lecture.pdf"
    output = tmp_path / "figures"
    _image_only_pdf(pdf, text="A labeled source figure.")

    chunks = ingest(pdf, extract_images=output)

    saved = output / "lecture-p1-img1.png"
    assert saved.is_file()
    assert saved.read_bytes()
    assert f"[figure: {saved} | lecture.pdf p.1]" in chunks[0].text


def test_extract_images_skips_sub_32px_image(tmp_path):
    pdf = tmp_path / "icon.pdf"
    output = tmp_path / "figures"
    _image_only_pdf(pdf, image_px=31, text="Text beside a decorative icon.")

    chunks = ingest(pdf, extract_images=output)

    assert list(output.iterdir()) == []
    assert "[figure:" not in chunks[0].text


def test_extract_images_default_none_regression(tmp_path):
    pdf = tmp_path / "default.pdf"
    output = tmp_path / "figures"
    _image_only_pdf(pdf, text="Default ingestion stays text-only.")

    chunks = ingest(pdf)

    assert not output.exists()
    assert "[figure:" not in chunks[0].text


def test_extract_images_image_only_page_emits_chunk(tmp_path):
    pdf = tmp_path / "scanned.pdf"
    output = tmp_path / "figures"
    _image_only_pdf(pdf)

    chunks = ingest(pdf, extract_images=output)

    saved = output / "scanned-p1-img1.png"
    assert saved.is_file()
    assert len(chunks) == 1
    assert "image-only page" in chunks[0].text
    assert f"[figure: {saved} | scanned.pdf p.1]" in chunks[0].text


def _bordered_table_pdf(path: Path) -> None:
    """Write a one-page PDF containing a 2x2 bordered table with cell text."""
    import fitz  # PyMuPDF

    doc = fitz.open()
    page = doc.new_page()
    cells = [["Metric", "Meaning"], ["ROI", "Return on investment"]]
    x0, y0, col_w, row_h = 72, 72, 160, 28
    for r in range(2):
        for c in range(2):
            rect = fitz.Rect(
                x0 + c * col_w, y0 + r * row_h,
                x0 + (c + 1) * col_w, y0 + (r + 1) * row_h,
            )
            page.draw_rect(rect, color=(0, 0, 0), width=0.7)
            page.insert_text((rect.x0 + 3, rect.y0 + 18), cells[r][c], fontsize=11)
    doc.save(str(path))
    doc.close()


def test_pdf_tables_are_re_emitted_as_structured_pipe_rows(tmp_path):
    pdf = tmp_path / "report.pdf"
    _bordered_table_pdf(pdf)

    chunks = ingest(pdf)

    assert len(chunks) == 1
    assert chunks[0].source == "report.pdf p.1"
    # The relational row/column structure is preserved alongside the flat text.
    assert "Table:" in chunks[0].text
    assert "Metric | Meaning" in chunks[0].text
    assert "ROI | Return on investment" in chunks[0].text


def test_pdf_prose_page_has_no_table_section(tmp_path):
    import fitz  # PyMuPDF

    pdf = tmp_path / "prose.pdf"
    doc = fitz.open()
    doc.new_page().insert_text((72, 72), "Photosynthesis converts light energy.")
    doc.save(str(pdf))
    doc.close()

    chunks = ingest(pdf)

    assert len(chunks) == 1
    assert "Photosynthesis" in chunks[0].text
    assert "Table:" not in chunks[0].text  # prose must not be mislabeled tabular


def test_pdf_math_like_text_layer_line_gets_provenance_label(tmp_path):
    import fitz  # PyMuPDF

    pdf = tmp_path / "equations.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Kinetic energy")
    page.insert_text((72, 96), "E = m * c^2")
    doc.save(str(pdf))
    doc.close()

    chunks = ingest(pdf)

    assert "E = m * c^2" in chunks[0].text
    assert "[math: E = m * c^2 | equations.pdf p.1 text layer]" in chunks[0].text


def test_pdf_math_marker_preserves_greek_and_operator_spans():
    class FakePage:
        def get_text(self, kind, sort):
            assert (kind, sort) == ("dict", True)
            return {
                "blocks": [
                    {
                        "type": 0,
                        "lines": [
                            {
                                "spans": [
                                    {"text": "\u0394E "},
                                    {"text": "= "},
                                    {"text": "\u03b1 + \u03b2"},
                                ]
                            }
                        ],
                    }
                ]
            }

    assert _extract_pdf_math(FakePage(), "physics.pdf p.2") == [
        "[math: \u0394E = \u03b1 + \u03b2 | physics.pdf p.2 text layer]"
    ]


def test_pdf_plain_prose_is_not_labeled_as_math(tmp_path):
    import fitz  # PyMuPDF

    pdf = tmp_path / "prose.pdf"
    doc = fitz.open()
    doc.new_page().insert_text((72, 72), "Revenue increased after the launch.")
    doc.save(str(pdf))
    doc.close()

    chunks = ingest(pdf)

    assert "[math:" not in chunks[0].text


def test_pptx_chart_preserves_categories_series_and_values(tmp_path):
    from pptx import Presentation
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    pptx = tmp_path / "chart.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    data = ChartData()
    data.categories = ["Q1", "Q2"]
    data.add_series("North", (10, 12.5))
    data.add_series("South", (8, 9))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(5),
        Inches(3),
        data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Quarterly revenue"
    prs.save(str(pptx))

    chunks = ingest(pptx)

    assert "Chart:\nTitle: Quarterly revenue" in chunks[0].text
    assert "Category | North | South" in chunks[0].text
    assert "Q1 | 10 | 8" in chunks[0].text
    assert "Q2 | 12.5 | 9" in chunks[0].text


def test_unknown_extension_raises(tmp_path):
    weird = tmp_path / "image.png"
    weird.write_bytes(b"\x89PNG")
    with pytest.raises(ValueError, match="unsupported"):
        ingest(weird)


def test_missing_file_raises(tmp_path):
    with pytest.raises(FileNotFoundError):
        ingest(tmp_path / "nope.md")


def test_cli_prints_provenance_and_text(tmp_path):
    md = tmp_path / "notes.md"
    md.write_text("Krebs cycle happens in the mitochondria.")

    proc = subprocess.run(
        [sys.executable, "scripts/ingest.py", str(md)],
        cwd=_REPO_ROOT, capture_output=True, text=True,
    )

    assert proc.returncode == 0, proc.stderr
    assert "notes.md" in proc.stdout
    assert "Krebs cycle" in proc.stdout
